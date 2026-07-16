# 파일 이름: pptx_converter.py
import re
import json
import os
from pptx import Presentation

def is_mostly_korean(text, threshold=0.4):
    """텍스트가 주로 한국어로 구성되어 있는지 확인합니다."""
    if not text:
        return False
    korean_chars = len(re.findall(r'[\uac00-\ud7a3]', text))
    return (korean_chars / len(text)) > threshold if len(text) > 0 else False

def convert_pptx_to_js(pptx_path):
    """
    PPTX 파일을 읽어 'opicData' JavaScript 변수 형식으로 변환합니다.
    복잡한 슬라이드 구조를 더 잘 처리하도록 개선된 버전입니다.
    """
    if not os.path.exists(pptx_path):
        error_message = f"// 오류: '{pptx_path}' 파일을 찾을 수 없습니다. 스크립트와 동일한 폴더에 있는지 확인해주세요."
        print(error_message)
        return error_message

    try:
        print(f"'{pptx_path}' 파일 분석을 시작합니다...")
        prs = Presentation(pptx_path)
    except Exception as e:
        error_message = f"// 오류: PPTX 파일을 여는 중 문제가 발생했습니다: {e}"
        print(error_message)
        return error_message

    all_topics = []
    current_topic = None

    for i, slide in enumerate(prs.slides):
        # 슬라이드에서 비어있지 않은 모든 텍스트 상자 추출
        slide_texts = [shape.text.strip() for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
        if not slide_texts:
            continue

        # --- 휴리스틱 1: 주제(Topic) 제목 슬라이드 식별 ---
        # "## 주제" 또는 "## Role-Play" 같은 키워드를 포함하는 단순한 구조의 슬라이드를 찾습니다.
        is_title_slide = False
        potential_title = ""
        if len(slide_texts) < 4: # 제목 슬라이드는 보통 텍스트 상자가 적습니다.
            for text in slide_texts:
                # "01 ...주제", "21 Role-Play..." 같은 제목을 찾는 정규식
                if re.match(r'^\d{1,2}\s*([\w\s&]+주제|Role-Play|ROLE PLAY)', text, re.IGNORECASE):
                    is_title_slide = True
                    potential_title = text.replace('\n', ' ').strip()
                    break
        
        if is_title_slide:
            # 새 주제를 시작하기 전에, 이전 주제에 내용이 있다면 저장합니다.
            if current_topic and current_topic["slides"]:
                all_topics.append(current_topic)
            
            current_topic = {
                "title": potential_title,
                "slides": []
            }
            print(f"\n  - 주제 발견 (슬라이드 {i+1}): {potential_title}")
            continue # 제목을 찾았으면 다음 슬라이드로 넘어갑니다.
        
        # --- 휴리스틱 2: 내용(Content) 슬라이드 처리 ---
        # 현재 주제에 속한 내용 슬라이드로 처리합니다.
        elif current_topic:
            korean_blocks = []
            english_blocks = []

            # 텍스트 상자를 한글과 영어로 분리합니다.
            for text in slide_texts:
                # 슬라이드 번호처럼 보이는 짧은 텍스트는 무시합니다.
                if text.isdigit() or len(text.strip()) < 20:
                    continue
                
                if is_mostly_korean(text):
                    korean_blocks.append(text)
                else:
                    english_blocks.append(text)

            # 한글과 영어 블록이 모두 존재할 경우에만 페어링을 시도합니다.
            if not (korean_blocks and english_blocks):
                if slide_texts: # 내용이 있는데도 블록이 비어있으면 경고
                    print(f"    - 경고: 슬라이드 {i+1}에서 한/영 스크립트 쌍을 찾지 못했습니다. 건너뜁니다.")
                continue

            # 휴리스틱: 슬라이드에 한/영 블록이 하나씩만 있는 가장 일반적인 경우
            if len(korean_blocks) == 1 and len(english_blocks) == 1:
                slide_data = {"korean": korean_blocks[0], "english": english_blocks[0]}
                current_topic["slides"].append(slide_data)
                print(f"    - 내용 추가 (1:1 페어링, 슬라이드 {i+1})")
                continue

            # 복잡한 경우: 슬라이드에 여러 스크립트가 섞여 있을 때 제목으로 페어링
            print(f"    - 복잡한 슬라이드 {i+1} 발견. 제목 기반 페어링 시도...")
            
            def get_script_title(text_block):
                # 첫 줄을 제목으로 간주하고 정규화합니다 (예: "02 [Adv] ... ♣" -> "02 [Adv]")
                first_line = text_block.split('\n')[0].strip()
                match = re.match(r'^\d{1,2}\s*\[(Int|Adv)\]', first_line)
                return match.group(0) if match else None

            korean_map = {get_script_title(b): b for b in korean_blocks if get_script_title(b)}
            english_map = {get_script_title(b): b for b in english_blocks if get_script_title(b)}

            matched_count = 0
            for title, ko_text in korean_map.items():
                if title in english_map:
                    en_text = english_map[title]
                    slide_data = {"korean": ko_text, "english": en_text}
                    current_topic["slides"].append(slide_data)
                    print(f"      - 성공: '{title}' 스크립트 페어링 완료")
                    del english_map[title] # 중복 매칭 방지
                    matched_count += 1
                else:
                    print(f"      - 실패: 한글 스크립트 '{title}'에 해당하는 영문 짝을 찾지 못함")
            
            if matched_count == 0:
                 print(f"    - 경고: 슬라이드 {i+1}에서 제목 기반 페어링에 실패했습니다.")


    # 마지막으로 처리된 주제를 리스트에 추가합니다.
    if current_topic and current_topic["slides"]:
        all_topics.append(current_topic)

    if not all_topics:
        return "// 오류: PPTX 파일에서 유효한 주제와 내용을 추출하지 못했습니다. 파일 구조를 확인해주세요."

    # 불필요한 HTML 필드를 제거하고 순수 텍스트 데이터만 생성합니다.
    clean_topics = [{"title": t["title"], "slides": [{"korean": s["korean"], "english": s["english"]} for s in t["slides"]]} for t in all_topics]

    js_string = json.dumps(clean_topics, indent=4, ensure_ascii=False)
    return f"const opicData = {js_string};"


# --- 스크립트 사용법 ---
# 1. 아래 파일 이름을 실제 변환하고 싶은 PPTX 파일 이름으로 수정하세요.
pptx_file_name = "rev05-2_영문_국문_NewOpic_답변_작은글씨.pptx"

# 2. 스크립트를 실행하면 변환된 코드가 'generated_opicData.js' 파일로 저장됩니다.
js_output = convert_pptx_to_js(pptx_file_name)

output_filename = "generated_opicData.js"
with open(output_filename, "w", encoding="utf-8") as f:
    f.write(js_output)

if "// 오류:" not in js_output:
    print(f"\n✅ 변환 성공! 결과가 '{output_filename}' 파일에 저장되었습니다.")
    print("   이 파일의 전체 내용을 복사하여 opic.html 파일의 기존 'const opicData = [...]' 부분을 교체하세요.")
else:
    print(f"\n❌ 변환 실패. '{output_filename}' 파일에서 오류 메시지를 확인하세요.")
